#!/usr/bin/env python3

import pandas as pd
import openpyxl as opxl
import shutil
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import RGBColor as docxColor
from datetime import date
from datetime import datetime
from dateutil.relativedelta import relativedelta

class company:
    """
    company object for the use of storing info gathered from the excel sheet
    """
    def __init__(self):
        self.fullName = ""
        self.companyName = ""
        self.email = ""
        self.phoneNumber = ""
        self.membershipType = ""
        self.space = ""
        self.description = ""
        self.term = ""
        self.startDate = ""
        self.endDate = ""
        self.rate = ""
        self.additionalFees = ""
        self.discountPromo = ""
        self.misc = ""
        self.deposit = ""
        self.link = ""
    
def gatherInfo(obj, sheetname):
    """
    Takes a comapny object and a specifit excel spreadsheet, and gathers the data from the excel and places it into the company object
    Note: This gathering only works for the specified Excel workbook that is used in this project
    """
    wb = opxl.load_workbook(sheetname)
    wb._active_sheet_index = 0
    ws = wb.active

    obj.fullName = ws["B1"].value
    obj.companyName = ws["B2"].value
    obj.email = ws["B3"].value
    obj.phoneNumber = ws["B4"].value

    wb._active_sheet_index = 1
    ws = wb.active

    obj.membershipType = ws["B1"].value
    obj.space = ws["B3"].value
    obj.description = ws["B4"].value
    obj.term = ws["B5"].value
    obj.startDate = ws["B6"].value
    obj.endDate = ws["B7"].value
    obj.rate = ws["B8"].value
    obj.additionalFees = ws["B9"].value
    obj.discountPromo = ws["B10"].value
    obj.misc = ws["B12"].value
    obj.deposit = ws["B13"].value
    obj.link = ws["B17"].value

def edit_ppt_text(pptFile, obj, path):
    """
    Takes the template pptx, company object, and current path/working directory and updates the template pptx with
    the data from the excel sheet and makes a copy of the template to store in the same directory
    """
    editDict = {}
    editDict["space1"] = obj.space
    editDict["description1"] = obj.description
    editDict["term1"] = obj.term
    editDict["price1"] = obj.rate


    presentation = Presentation(pptFile)

    for item in presentation.slides[0].shapes:
        if hasattr(item, "text"):
            if item.text == "Company name":
                item.text = obj.companyName + " Proposal"      
                for paragraph in item.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Bebas Neue"
                        run.font.size = Pt(72)
                        run.font.color.rgb = RGBColor(255,255,255)

    for item in presentation.slides[5].shapes:
        if hasattr(item, "text"):
            if item.text == "Quote for XYZ Company":
                item.text = "Quote for " + obj.companyName + " Company"
                for paragraph in item.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Bebas Neue"
                        run.font.size = Pt(66)
                        run.font.color.rgb = RGBColor(0,0,0)
            #TODO: Once John sends me specs on desc and others, I can start working on auto-formatting
            if item.text in editDict:
                item.text = str(editDict[item.text])
                for paragraph in item.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.name = "Open Sans"
                        run.font.size = Pt(22)
                        run.font.color.rgb = RGBColor(0,0,0)

    presentation.save(f"{path}/{obj.companyName} VX Proposal.pptx")   

def replace_text(doc, old_text, new_text):
    """
    Replace old_text with new_text in all paragraphs of the docx document.
    """
    for paragraph in doc.paragraphs:
        if old_text in paragraph.text:
            for run in paragraph.runs:
                run.text = run.text.replace(old_text, new_text)
                run.font.color.rgb = docxColor(0, 0, 0)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if old_text in cell.text:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.replace(old_text, new_text)
                            run.font.color.rgb = docxColor(0, 0, 0)

def edit_docx_text(docxFile, obj, path):
    """
    Goes through the whole document and replaces all items that need to be replaced
    """
    #Handles date information
    today = datetime.today()
    today = today.strftime("%m/%d/%Y")
    currentDate = datetime.now()

    #Adds term to current date, term need to be in months
    endTermDate = currentDate + relativedelta(months=obj.term)
    endTermDate = endTermDate.strftime("%m/%d/%Y")

    doc = Document(docxFile)

    replace_text(doc, "DATEREPLACE", today)
    replace_text(doc, "DATEREPLACE", today)
    replace_text(doc, "DATEREPLACE", today)
    replace_text(doc, "NAMEREPLACE", obj.companyName)
    replace_text(doc, "SPACEREPLACE", obj.space)
    replace_text(doc, "TERMREPLACE", str(obj.term))
    replace_text(doc, "EXPIRATIONREPLACE", endTermDate)
    #TODO: Fix below, need info from John
    replace_text(doc, "DEPOSITREPLACE", f"${obj.deposit}")

    doc.save(f"{path}/{obj.companyName} VX Membership Agreement.docx")        

#Gets the directory location of the script
script_dir = os.path.dirname(os.path.realpath(__file__))
#Changes current directory to that of the script
os.chdir(script_dir)
#Saves the path to the script in path
path = os.getcwd()
#Gets the template pptx
sourcePPT = f"{path}/VX Proposal Template.pptx"
#Gets the template docx
sourceDoc = f"{path}/Venture X Membership Agreement And T&C Generic Jan 24.docx"
#Gets the data excel sheet
sheetname = f"{path}/Membership Proposal Details.xlsx"

#Creates an instance of company class and populates its variables with gather info
obj = company()
gatherInfo(obj, sheetname)

#Calls the edit_ppt_text method
edit_ppt_text(sourcePPT, obj, path)

#Calls the edit_docx_text method
edit_docx_text(sourceDoc, obj, path)